from __future__ import annotations

import asyncio
import os
import io
import json
import logging
import tempfile
import traceback
import subprocess
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Optional
from urllib.parse import urlparse

import PyPDF2
from fastapi import APIRouter, Body, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, Response
from pydantic import BaseModel, Field
from sqlalchemy.orm import Session

import models
from database import get_db
from deps import call_ai, call_ai_async, get_current_user, get_user_by_email, get_user_by_username, optional_security, SECRET_KEY, ALGORITHM, JWT_AUDIENCE
from services.api_key_pool import ApiKeyPoolExhausted, is_provider_quota_error, provider_limit_exhausted
from services.document_processor import extract_text_from_pdf_detailed
from services.storage_service import StorageService

try:
    from services.ai_media_processor import ai_media_processor
except ImportError:
    ai_media_processor = None

try:
    from services.math_processor import process_math_in_response
except ImportError:
    process_math_in_response = None

try:
    from services.podcast_agent import podcast_agent_service
except ImportError:
    podcast_agent_service = None

logger = logging.getLogger(__name__)

UPLOAD_DIR = Path("uploads/slides")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
SLIDE_CACHE_DIR = Path(os.getenv("SLIDE_CACHE_DIR", "uploads/slide_cache"))
SLIDE_CACHE_DIR.mkdir(parents=True, exist_ok=True)
MEDIA_UPLOAD_MAX_MB = max(1, int(os.getenv("MEDIA_UPLOAD_MAX_MB", "100")))
MEDIA_UPLOAD_MAX_BYTES = MEDIA_UPLOAD_MAX_MB * 1024 * 1024

router = APIRouter(prefix="/api", tags=["media"])

def _is_remote_storage_path(file_path: str | None) -> bool:
    return bool(file_path and urlparse(file_path).scheme in {"s3", "r2"})

def _storage_key_from_uri(file_path: str) -> str:
    parsed = urlparse(file_path)
    if parsed.scheme in {"s3", "r2"}:
        return parsed.path.lstrip("/")
    return file_path

def _local_slide_cache_path(slide) -> Path:
    safe_name = "".join(c for c in (slide.filename or f"slide_{slide.id}") if c.isalnum() or c in (".", "_", "-"))
    return SLIDE_CACHE_DIR / f"{slide.id}_{safe_name}"

def _ensure_local_slide_file(slide) -> Path:
    file_path = slide.file_path or ""
    if not _is_remote_storage_path(file_path):
        local_path = Path(file_path)
        if not local_path.exists():
            raise HTTPException(status_code=404, detail="Slide file not found")
        return local_path

    cache_path = _local_slide_cache_path(slide)
    if cache_path.exists() and cache_path.stat().st_size > 0:
        return cache_path

    storage = StorageService.get_storage()
    if getattr(storage, "storage_type", "local") == "local":
        raise HTTPException(status_code=500, detail="Remote storage is not configured")
    storage.download_file(_storage_key_from_uri(file_path), cache_path)
    return cache_path

class PodcastStartRequest(BaseModel):
    user_id: str
    transcript: str
    analysis: dict = Field(default_factory=dict)
    title: str = "Media Podcast"
    source_type: str = "media"
    voice_mode: str = "coach"
    voice_persona: str = "mentor"
    difficulty: str = "intermediate"
    answer_language: str = "en"
    session_options: dict = Field(default_factory=dict)

class PodcastSessionRequest(BaseModel):
    user_id: str
    session_id: str

class PodcastQuestionRequest(BaseModel):
    user_id: str
    session_id: str
    question: str
    voice_mode: Optional[str] = None
    voice_persona: Optional[str] = None
    difficulty: Optional[str] = None
    question_language: Optional[str] = None
    answer_language: Optional[str] = None

class PodcastVoiceModeRequest(BaseModel):
    user_id: str
    session_id: str
    voice_mode: str

class PodcastSettingsRequest(BaseModel):
    user_id: str
    session_id: str
    voice_mode: Optional[str] = None
    voice_persona: Optional[str] = None
    difficulty: Optional[str] = None
    answer_language: Optional[str] = None
    session_options: dict = Field(default_factory=dict)

class PodcastBookmarkRequest(BaseModel):
    user_id: str
    session_id: str
    label: Optional[str] = None
    chapter_index: Optional[int] = None
    timestamp_seconds: Optional[int] = None

class PodcastReplayBookmarkRequest(BaseModel):
    user_id: str
    session_id: str
    bookmark_id: int

class PodcastJumpChapterRequest(BaseModel):
    user_id: str
    session_id: str
    chapter_index: int

class PodcastMCQStartRequest(BaseModel):
    user_id: str
    session_id: str
    count: int = 5

class PodcastMCQAnswerRequest(BaseModel):
    user_id: str
    session_id: str
    question_index: int
    selected_index: int

class PodcastExportRequest(BaseModel):
    user_id: str
    session_id: str
    format_type: str = "markdown"

def _resolve_user_or_404(db: Session, user_id: str):
    user = get_user_by_username(db, user_id) or get_user_by_email(db, user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    return user

def _verify_user_access(requested_user, current_user: models.User):
    if requested_user.id != current_user.id:
        raise HTTPException(status_code=403, detail="Access denied")

def _safe_audio_suffix(filename: Optional[str]) -> str:
    suffix = Path(filename or "").suffix.lower()
    return suffix if suffix in {".m4a", ".mp3", ".wav", ".webm", ".ogg", ".opus", ".mp4", ".mpeg", ".mpga"} else ".webm"

def _is_pdf_upload(file: UploadFile, content: bytes) -> bool:
    filename = (file.filename or "").lower()
    content_type = (file.content_type or "").lower()
    return (
        filename.endswith(".pdf")
        or content_type == "application/pdf"
        or content.startswith(b"%PDF")
    )

def _safe_upload_name(filename: Optional[str]) -> str:
    return "".join(c for c in (filename or "upload") if c.isalnum() or c in (".", "_", "-")) or "upload"

def _storage_uri_for_path(storage, storage_path: str) -> str:
    if getattr(storage, "storage_type", "local") != "local" and hasattr(storage, "uri_for_path"):
        return storage.uri_for_path(storage_path)
    return storage_path

def _store_media_source(storage, content: bytes, *, user_id: int, filename: Optional[str], content_type: Optional[str]) -> str:
    safe_name = _safe_upload_name(filename)
    timestamp = int(datetime.now(timezone.utc).timestamp())
    if getattr(storage, "storage_type", "local") == "local":
        upload_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "uploads", "media")
        os.makedirs(upload_dir, exist_ok=True)
        file_path = os.path.join(upload_dir, f"{user_id}_{timestamp}_{safe_name}")
        with open(file_path, "wb") as f:
            f.write(content)
        return file_path

    storage_key = f"media_uploads/{user_id}/{timestamp}_{safe_name}"
    result = storage.upload_bytes(content, storage_key, content_type or "application/octet-stream")
    stored_path = result.get("storage_path") or storage_key
    return _storage_uri_for_path(storage, stored_path)

def _read_media_source_bytes(storage, file_path: str, fallback: bytes) -> bytes:
    if getattr(storage, "storage_type", "local") == "local":
        return fallback
    if not hasattr(storage, "download_bytes"):
        raise HTTPException(status_code=500, detail="Configured storage does not support file readback")
    try:
        return storage.download_bytes(_storage_key_from_uri(file_path))
    except Exception as exc:
        logger.error("Failed to read media upload from storage: %s", exc, exc_info=True)
        raise HTTPException(status_code=502, detail="Uploaded media was stored but could not be read back from S3")

@router.post("/transcribe_audio/")
async def transcribe_audio(
    audio_file: UploadFile = File(...),
    user_id: str = Form(...),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    temp_audio_path = None
    try:
        requested_user = _resolve_user_or_404(db, user_id)
        _verify_user_access(requested_user, current_user)
        user = requested_user

        audio_content = await audio_file.read()
        if len(audio_content) == 0:
            raise HTTPException(status_code=400, detail="Empty audio file")

        with tempfile.NamedTemporaryFile(delete=False, suffix=_safe_audio_suffix(audio_file.filename), mode="wb") as temp_file:
            temp_file.write(audio_content)
            temp_audio_path = temp_file.name

        try:
            from groq import Groq
            groq_api_key = os.getenv("GROQ_API_KEY")
            groq_client = Groq(api_key=groq_api_key) if groq_api_key else None
            if not groq_client:
                raise HTTPException(status_code=500, detail="Transcription service not configured")

            with open(temp_audio_path, "rb") as f:
                transcription = groq_client.audio.transcriptions.create(
                    file=f,
                    model="whisper-large-v3-turbo",
                    response_format="json",
                    language="en",
                )

            transcript_text = transcription.text
            return {
                "status": "success",
                "transcript": transcript_text,
                "length": len(transcript_text),
                "model_used": "whisper-large-v3-turbo",
            }

        except HTTPException:
            raise
        except Exception as groq_error:
            if is_provider_quota_error(groq_error):
                raise provider_limit_exhausted("groq") from groq_error
            logger.error(f"Groq API error: {str(groq_error)}")
            raise HTTPException(status_code=500, detail="Transcription failed")

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error in transcribe_audio: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to transcribe audio")
    finally:
        if temp_audio_path and os.path.exists(temp_audio_path):
            try:
                os.remove(temp_audio_path)
            except Exception as cleanup_error:
                logger.warning(f"Failed to cleanup temp file: {cleanup_error}")

@router.get("/test_transcribe")
def test_transcribe_endpoint():
    groq_api_key = os.getenv("GROQ_API_KEY")
    return {
        "status": "endpoint exists",
        "message": "Transcribe audio endpoint is registered",
        "groq_available": groq_api_key is not None,
    }

@router.post("/transcribe_audio_test/")
async def transcribe_audio_test(
    user_id: str = Form(...),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    try:
        requested_user = _resolve_user_or_404(db, user_id)
        _verify_user_access(requested_user, current_user)
        user = requested_user

        groq_api_key = os.getenv("GROQ_API_KEY")
        return {
            "status": "success",
            "message": "Endpoint working, user found",
            "user_id": user.id,
            "groq_configured": groq_api_key is not None,
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Test endpoint error: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal server error")

@router.post("/media/process")
async def process_media(
    user_id: str = Form(...),
    file: UploadFile = File(None),
    youtube_url: str = Form(None),
    note_style: str = Form("detailed"),
    difficulty: str = Form("intermediate"),
    subject: str = Form("general"),
    custom_instructions: str = Form(None),
    generate_flashcards: bool = Form(False),
    generate_quiz: bool = Form(False),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    try:
        requested_user = _resolve_user_or_404(db, user_id)
        _verify_user_access(requested_user, current_user)
        user = requested_user

        if not ai_media_processor:
            raise HTTPException(status_code=500, detail="Media processor not available")

        transcript_data = None
        analysis_data = None
        source_type = None
        filename = None
        file_path = None
        source_storage_type = None

        logger.info(f"Received - File: {file.filename if file else None}, YouTube URL: '{youtube_url}'")

        if youtube_url and youtube_url.strip():
            logger.info(f"Processing YouTube URL: {youtube_url}")
            result = await ai_media_processor.process_youtube_video(youtube_url.strip())

            if not result.get("success"):
                raise HTTPException(status_code=400, detail=result.get("error", "Failed to process YouTube video"))

            transcript_data = result
            source_type = "youtube"
            filename = result["video_info"]["title"]

        elif file:
            logger.info(f"Processing uploaded file: {file.filename}")
            content = await file.read()
            if not content:
                raise HTTPException(status_code=400, detail="The uploaded file is empty")
            if len(content) > MEDIA_UPLOAD_MAX_BYTES:
                raise HTTPException(
                    status_code=413,
                    detail=f"File is too large. Maximum upload size is {MEDIA_UPLOAD_MAX_MB}MB",
                )

            storage = StorageService.get_storage()
            source_storage_type = getattr(storage, "storage_type", "local")
            file_path = _store_media_source(
                storage,
                content,
                user_id=user.id,
                filename=file.filename,
                content_type=file.content_type,
            )

            if _is_pdf_upload(file, content):
                logger.info("Extracting text from uploaded PDF")
                source_bytes = _read_media_source_bytes(storage, file_path, content)
                extraction = await asyncio.to_thread(extract_text_from_pdf_detailed, source_bytes)
                extracted_text = (extraction.get("text") or "").strip()
                if not extracted_text:
                    warnings = extraction.get("warnings", []) or []
                    needs_ocr = any("ocr" in str(w).lower() or "tesseract" in str(w).lower() for w in warnings)
                    raise HTTPException(
                        status_code=400,
                        detail=(
                            "No readable text was found in this scanned or image-only PDF, "
                            "and OCR is not available on the server. Install Tesseract OCR "
                            "and the pytesseract Python package, then retry."
                            if needs_ocr
                            else (
                                "No readable text was found in this PDF. "
                                "It may be scanned, image-only, encrypted, or malformed."
                            )
                        ),
                    )
                transcript_data = {
                    "success": True,
                    "transcript": extracted_text,
                    "language": "en",
                    "duration": 0,
                    "has_timestamps": False,
                    "segments": [],
                    "pdf_info": {
                        "page_count": extraction.get("page_count", 0),
                        "parser": extraction.get("parser", ""),
                        "warnings": extraction.get("warnings", []),
                    },
                }
                source_type = "pdf"
            else:
                transcription_path = file_path
                temp_media_path = None
                if getattr(storage, "storage_type", "local") != "local":
                    with tempfile.NamedTemporaryFile(delete=False, suffix=_safe_audio_suffix(file.filename), mode="wb") as temp_file:
                        temp_file.write(_read_media_source_bytes(storage, file_path, content))
                        temp_media_path = temp_file.name
                    transcription_path = temp_media_path

                try:
                    result = await ai_media_processor.transcribe_audio_groq(transcription_path)
                finally:
                    if temp_media_path and os.path.exists(temp_media_path):
                        try:
                            os.remove(temp_media_path)
                        except Exception as cleanup_error:
                            logger.warning(f"Failed to cleanup temp media file: {cleanup_error}")

                if not result.get("success"):
                    raise HTTPException(status_code=400, detail=result.get("error", "Failed to transcribe audio"))

                transcript_data = result
                source_type = "upload"
            filename = file.filename

        else:
            raise HTTPException(status_code=400, detail="Please provide either a file or YouTube URL")

        logger.info("Performing AI analysis...")
        analysis_result = await ai_media_processor.analyze_transcript_ai(
            transcript_data["transcript"],
            {"subject": subject, "difficulty": difficulty},
            user_id=user.id,
        )

        if not analysis_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"AI analysis failed: {analysis_result.get('error', 'Unknown error')}",
            )

        analysis_data = analysis_result.get("analysis", {})

        logger.info(f"Generating {note_style} notes...")
        notes_result = await ai_media_processor.generate_notes_ai(
            transcript_data["transcript"],
            analysis_data,
            note_style,
            {
                "difficulty": difficulty,
                "subject": subject,
                "custom_instructions": custom_instructions,
            },
            user_id=user.id,
        )

        if not notes_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"Note generation failed: {notes_result.get('error', 'Unknown error')}",
            )

        flashcards = []
        if generate_flashcards:
            logger.info("Generating flashcards...")
            flashcard_result = await ai_media_processor.generate_flashcards_ai(
                transcript_data["transcript"],
                analysis_data,
                count=10,
                user_id=user.id,
            )
            if flashcard_result.get("success"):
                flashcards = flashcard_result.get("flashcards", [])
            else:
                logger.warning(f"Flashcard generation failed: {flashcard_result.get('error', 'Unknown')}")

        quiz_questions = []
        if generate_quiz:
            logger.info("Generating quiz...")
            quiz_result = await ai_media_processor.generate_quiz_ai(
                transcript_data["transcript"],
                analysis_data,
                count=10,
                user_id=user.id,
            )
            if quiz_result.get("success"):
                quiz_questions = quiz_result.get("questions", [])
            else:
                logger.warning(f"Quiz generation failed: {quiz_result.get('error', 'Unknown')}")

        key_moments = []
        if transcript_data.get("has_timestamps") and transcript_data.get("segments"):
            key_moments = await ai_media_processor.extract_key_moments(
                transcript_data["segments"],
                analysis_data,
            )

        response = {
            "success": True,
            "source_type": source_type,
            "filename": filename,
            "transcript": transcript_data["transcript"],
            "language": transcript_data.get("language", "en"),
            "language_name": ai_media_processor.get_language_name(transcript_data.get("language", "en")),
            "duration": transcript_data.get("duration", 0),
            "has_timestamps": transcript_data.get("has_timestamps", False),
            "segments": transcript_data.get("segments", [])[:50],
            "pdf_info": transcript_data.get("pdf_info") if source_type == "pdf" else None,
            "source_storage_path": file_path,
            "source_storage_type": source_storage_type if file_path else None,
            "analysis": analysis_data,
            "notes": {
                "content": notes_result["content"],
                "style": note_style,
            },
            "flashcards": flashcards,
            "quiz_questions": quiz_questions,
            "key_moments": key_moments,
            "video_info": transcript_data.get("video_info") if source_type == "youtube" else None,
        }

        return response

    except HTTPException:
        raise
    except ApiKeyPoolExhausted:
        raise
    except Exception as e:
        logger.error(f"Media processing error: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail="Processing failed")

@router.post("/media/save-notes")
async def save_media_notes(
    user_id: str = Body(...),
    title: str = Body(...),
    content: str = Body(...),
    transcript: str = Body(None),
    analysis: dict = Body(None),
    flashcards: list = Body(None),
    quiz_questions: list = Body(None),
    key_moments: list = Body(None),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    try:
        requested_user = _resolve_user_or_404(db, user_id)
        _verify_user_access(requested_user, current_user)
        user = requested_user

        new_note = models.Note(
            user_id=user.id,
            title=title,
            content=content,
            custom_font="__MEDIA_NOTE__",
            transcript=transcript,
            analysis=json.dumps(analysis) if analysis else None,
            flashcards=json.dumps(flashcards) if flashcards else None,
            quiz_questions=json.dumps(quiz_questions) if quiz_questions else None,
            key_moments=json.dumps(key_moments) if key_moments else None,
            created_at=datetime.now(timezone.utc),
            updated_at=datetime.now(timezone.utc),
        )

        db.add(new_note)
        db.commit()
        db.refresh(new_note)

        if flashcards and len(flashcards) > 0:
            flashcard_set = models.FlashcardSet(
                user_id=user.id,
                title=f"Flashcards: {title}",
                source_type="media",
                source_id=new_note.id,
            )
            db.add(flashcard_set)
            db.commit()
            db.refresh(flashcard_set)

            for fc in flashcards:
                flashcard = models.Flashcard(
                    set_id=flashcard_set.id,
                    question=fc.get("question", ""),
                    answer=fc.get("answer", ""),
                    difficulty=fc.get("difficulty", "medium"),
                )
                db.add(flashcard)

            db.commit()

        from services.gamification_system import award_points
        award_points(db, user.id, "note_created")

        return {
            "success": True,
            "note_id": new_note.id,
            "message": "Notes saved successfully",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error saving notes: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to save notes")

@router.post("/media/regenerate-notes")
async def regenerate_notes(
    transcript: str = Body(...),
    analysis: dict = Body(...),
    note_style: str = Body("detailed"),
    difficulty: str = Body("intermediate"),
    subject: str = Body("general"),
    custom_instructions: str = Body(None),
    current_user: models.User = Depends(get_current_user),
):
    try:
        if not ai_media_processor:
            raise HTTPException(status_code=500, detail="Media processor not available")

        notes_result = await ai_media_processor.generate_notes_ai(
            transcript,
            analysis,
            note_style,
            {
                "difficulty": difficulty,
                "subject": subject,
                "custom_instructions": custom_instructions,
            },
            user_id=current_user.id,
        )

        if not notes_result.get("success"):
            raise HTTPException(status_code=500, detail="Note regeneration failed")

        return {
            "success": True,
            "content": notes_result["content"],
            "style": note_style,
        }

    except HTTPException:
        raise
    except ApiKeyPoolExhausted:
        raise
    except Exception as e:
        logger.error(f"Error regenerating notes: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to regenerate notes")

@router.get("/media/estimate-cost")
async def estimate_processing_cost(
    duration_seconds: int = Query(...),
    file_size_mb: float = Query(...),
):
    if not ai_media_processor:
        raise HTTPException(status_code=500, detail="Media processor not available")
    cost_info = ai_media_processor.estimate_processing_cost(duration_seconds, file_size_mb)
    return cost_info

@router.post("/media/generate-title")
async def generate_smart_title(
    transcript: str = Body(...),
    key_concepts: list = Body([]),
    summary: str = Body(""),
):
    try:
        prompt = f"""Generate a concise 3-4 word title for these notes.

Summary: {summary[:200]}
Key Concepts: {', '.join(key_concepts[:5])}
Content: {transcript[:500]}

Return ONLY the title, nothing else. Make it descriptive and catchy."""

        title = await call_ai_async(prompt, max_tokens=20, temperature=0.7)
        title = title.strip().strip('"').strip("'")
        words = title.split()
        if len(words) > 4:
            title = " ".join(words[:4])

        return {"title": title}

    except Exception as e:
        logger.error(f"Title generation error: {str(e)}")
        return {"title": "Media Notes"}

@router.get("/media/note/{note_id}")
async def get_single_note(
    note_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    try:
        note = db.query(models.Note).filter(
            models.Note.id == note_id,
            models.Note.user_id == current_user.id,
            models.Note.is_deleted == False,
        ).first()

        if not note:
            raise HTTPException(status_code=404, detail="Note not found")

        analysis = {}
        flashcards = []
        quiz_questions = []
        key_moments = []

        try:
            if note.analysis:
                analysis = json.loads(note.analysis)
        except Exception:
            pass

        try:
            if note.flashcards:
                flashcards = json.loads(note.flashcards)
        except Exception:
            pass

        try:
            if note.quiz_questions:
                quiz_questions = json.loads(note.quiz_questions)
        except Exception:
            pass

        try:
            if note.key_moments:
                key_moments = json.loads(note.key_moments)
        except Exception:
            pass

        return {
            "id": note.id,
            "title": note.title,
            "content": note.content,
            "created_at": note.created_at.isoformat(),
            "updated_at": note.updated_at.isoformat(),
            "folder_id": note.folder_id,
            "is_favorite": note.is_favorite,
            "custom_font": note.custom_font,
            "transcript": note.transcript or "",
            "analysis": analysis,
            "flashcards": flashcards,
            "quiz_questions": quiz_questions,
            "key_moments": key_moments,
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error fetching note {note_id}: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal server error")

@router.get("/media/history")
async def get_media_history(
    user_id: str = Query(...),
    limit: int = Query(50, ge=1, le=100),
    scan_limit: int = Query(50, ge=1, le=200),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    try:
        requested_user = _resolve_user_or_404(db, user_id)
        _verify_user_access(requested_user, current_user)
        user = requested_user

        all_notes = (
            db.query(models.Note)
            .filter(
                models.Note.user_id == user.id,
                models.Note.is_deleted == False,
            )
            .order_by(models.Note.created_at.desc())
            .limit(scan_limit)
            .all()
        )

        history = []
        for note in all_notes:
            is_media_note = note.custom_font == "__MEDIA_NOTE__"

            if not is_media_note:
                flashcard_set = (
                    db.query(models.FlashcardSet)
                    .filter(
                        models.FlashcardSet.source_type == "media",
                        models.FlashcardSet.source_id == note.id,
                    )
                    .first()
                )
                if flashcard_set:
                    is_media_note = True

            if is_media_note:
                history.append({
                    "id": note.id,
                    "title": note.title,
                    "created_at": note.created_at.isoformat(),
                    "preview": note.content[:200] if note.content else "",
                })
                if len(history) >= limit:
                    break

        logger.info(f"Found {len(history)} media-generated notes for user {user_id}")
        return {"history": history}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"History fetch error: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal server error")

_MAX_SLIDE_SIZE = 50 * 1024 * 1024
_ALLOWED_SLIDE_EXTENSIONS = {".pdf", ".ppt", ".pptx"}
_SLIDE_MAGIC_BYTES = {
    b'%PDF': 'pdf',
    b'PK\x03\x04': 'pptx',
    b'\xd0\xcf\x11\xe0': 'ppt',
}

@router.post("/upload_slides")
async def upload_slides(
    files: List[UploadFile] = File(...),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    try:
        user = current_user

        uploaded_slides = []

        for file in files:
            ext = Path(file.filename or "").suffix.lower()
            if ext not in _ALLOWED_SLIDE_EXTENSIONS:
                continue

            file_content = await file.read()
            file_size = len(file_content)

            if file_size > _MAX_SLIDE_SIZE:
                raise HTTPException(status_code=413, detail=f"File {file.filename} exceeds 50 MB limit")

            timestamp = int(datetime.now(timezone.utc).timestamp())
            clean_name = "".join(c for c in (file.filename or "upload") if c.isalnum() or c in (".", "_", "-"))
            safe_filename = f"{user.id}_{timestamp}_{clean_name}"
            local_file_path = UPLOAD_DIR / safe_filename

            with open(local_file_path, "wb") as f:
                f.write(file_content)

            storage = StorageService.get_storage()
            stored_file_path = str(local_file_path)
            if getattr(storage, "storage_type", "local") != "local":
                storage_key = f"slides/{user.id}/{safe_filename}"
                storage.upload_bytes(file_content, storage_key, file.content_type or "application/octet-stream")
                stored_file_path = (
                    storage.uri_for_path(storage_key)
                    if hasattr(storage, "uri_for_path")
                    else storage_key
                )

            page_count = 0
            extracted_text = ""

            if file.filename.lower().endswith(".pdf"):
                try:
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                    page_count = len(pdf_reader.pages)

                    for page in pdf_reader.pages[:10]:
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            extracted_text += page_text + "\n"

                    extracted_text = extracted_text[:10000]
                except Exception as e:
                    logger.error(f"Error extracting PDF text: {str(e)}")

            elif file.filename.lower().endswith((".ppt", ".pptx")):
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(file_content))
                    page_count = len(prs.slides)

                    for slide_idx, ppt_slide in enumerate(prs.slides):
                        if slide_idx >= 10:
                            break
                        for shape in ppt_slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                extracted_text += shape.text.strip() + "\n"

                    extracted_text = extracted_text[:10000]
                except Exception as e:
                    logger.error(f"Error extracting PowerPoint content: {str(e)}")

            slide = models.UploadedSlide(
                user_id=user.id,
                filename=safe_filename,
                original_filename=file.filename,
                file_path=stored_file_path,
                file_size=file_size,
                file_type=file.content_type or "application/pdf",
                page_count=page_count,
                extracted_text=extracted_text,
                processing_status="completed",
                uploaded_at=datetime.now(timezone.utc),
                processed_at=datetime.now(timezone.utc),
            )

            db.add(slide)
            uploaded_slides.append(slide)

        db.commit()

        return {
            "status": "success",
            "uploaded_count": len(uploaded_slides),
            "message": f"Successfully uploaded {len(uploaded_slides)} slide(s)",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading slides: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail="Failed to upload slides")

@router.get("/get_uploaded_slides")
def get_uploaded_slides(
    user_id: str = Query(...),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    try:
        requested_user = _resolve_user_or_404(db, user_id)
        _verify_user_access(requested_user, current_user)
        user = requested_user

        slides = (
            db.query(models.UploadedSlide)
            .filter(models.UploadedSlide.user_id == user.id)
            .order_by(models.UploadedSlide.uploaded_at.desc())
            .all()
        )

        return {
            "slides": [
                {
                    "id": slide.id,
                    "title": slide.original_filename,
                    "filename": slide.original_filename,
                    "file_size": slide.file_size,
                    "file_type": slide.file_type,
                    "page_count": slide.page_count,
                    "created_at": slide.uploaded_at.isoformat() + "Z",
                    "uploaded_at": slide.uploaded_at.isoformat() + "Z",
                    "preview_url": slide.preview_url,
                    "processing_status": slide.processing_status,
                }
                for slide in slides
            ]
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting slides: {str(e)}")
        return {"slides": []}

@router.delete("/delete_slide/{slide_id}")
def delete_slide(
    slide_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id,
            models.UploadedSlide.user_id == current_user.id,
        ).first()

        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

        stored_file_path = slide.file_path
        file_path = Path(stored_file_path).resolve() if stored_file_path and not _is_remote_storage_path(stored_file_path) else None
        allowed_dir = Path("uploads/slides").resolve()

        db.query(models.SlideAnalysis).filter(
            models.SlideAnalysis.slide_id == slide_id,
        ).delete(synchronize_session=False)
        db.query(models.QuestionSetSlide).filter(
            models.QuestionSetSlide.slide_id == slide_id,
        ).delete(synchronize_session=False)
        db.query(models.LearningReviewSlide).filter(
            models.LearningReviewSlide.slide_id == slide_id,
        ).delete(synchronize_session=False)
        db.delete(slide)
        db.commit()

        try:
            if _is_remote_storage_path(stored_file_path):
                storage = StorageService.get_storage()
                if getattr(storage, "storage_type", "local") != "local":
                    storage.delete_file(_storage_key_from_uri(stored_file_path))
                cache_path = _local_slide_cache_path(slide)
                if cache_path.exists():
                    cache_path.unlink()
            elif file_path and str(file_path).startswith(str(allowed_dir)) and file_path.exists():
                file_path.unlink()
        except Exception as cleanup_error:
            logger.warning(f"Deleted slide {slide_id}, but failed to remove file: {cleanup_error}")

        return {
            "status": "success",
            "message": "Slide deleted successfully",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting slide: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail="Failed to delete slide")

@router.get("/analyze_slide/{slide_id}")
async def analyze_slide(
    slide_id: int,
    force_reanalyze: bool = Query(False),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    try:
        from services.comprehensive_slide_analyzer import get_or_create_analysis

        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id,
            models.UploadedSlide.user_id == current_user.id,
        ).first()

        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

        file_path = _ensure_local_slide_file(slide)

        file_ext = slide.original_filename.lower().split(".")[-1]
        if file_ext not in ["pdf", "ppt", "pptx"]:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        analysis_result = get_or_create_analysis(
            slide_id=slide_id,
            file_path=file_path,
            file_type=file_ext,
            db=db,
            force_reanalyze=force_reanalyze,
        )

        return {
            "status": "success",
            "filename": slide.original_filename,
            "total_slides": analysis_result["total_slides"],
            "presentation_summary": analysis_result["presentation_summary"],
            "slides": analysis_result["slides"],
            "analyzed_at": analysis_result["analyzed_at"],
        }

    except HTTPException:
        raise
    except ApiKeyPoolExhausted:
        raise
    except Exception as e:
        logger.error(f"Error analyzing slide: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Error analyzing slide")

@router.get("/slide_file/{slide_id}")
async def get_slide_file(
    slide_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id,
            models.UploadedSlide.user_id == current_user.id,
        ).first()

        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

        file_path = _ensure_local_slide_file(slide)

        filename = slide.original_filename.lower()
        if filename.endswith(".pdf"):
            media_type = "application/pdf"
        elif filename.endswith(".pptx"):
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif filename.endswith(".ppt"):
            media_type = "application/vnd.ms-powerpoint"
        else:
            media_type = "application/octet-stream"

        return FileResponse(
            path=str(file_path),
            media_type=media_type,
            filename=slide.original_filename,
            headers={
                "Content-Disposition": f'inline; filename="{slide.original_filename}"'
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error serving slide file: {str(e)}")
        raise HTTPException(status_code=500, detail="Error serving slide file")

@router.get("/slide_image/{slide_id}/{page_number}")
async def get_slide_image(
    slide_id: int,
    page_number: int,
    token: Optional[str] = Query(default=None),
    credentials=Depends(optional_security),
    db: Session = Depends(get_db),
):
    from jose import JWTError, jwt as jose_jwt
    from fastapi.security import HTTPAuthorizationCredentials
    jwt_str = token or (credentials.credentials if credentials else None)
    if not jwt_str:
        raise HTTPException(status_code=401, detail="Authentication required")
    try:
        payload = jose_jwt.decode(jwt_str, SECRET_KEY, algorithms=[ALGORITHM], audience=JWT_AUDIENCE)
        username = payload.get("sub")
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid token")
    current_user = db.query(models.User).filter(models.User.username == username).first()
    if not current_user:
        current_user = db.query(models.User).filter(models.User.email == username).first()
    if not current_user:
        raise HTTPException(status_code=404, detail="User not found")
    import base64

    try:
        slide = db.query(models.UploadedSlide).filter(
            models.UploadedSlide.id == slide_id,
            models.UploadedSlide.user_id == current_user.id,
        ).first()

        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

        file_path = _ensure_local_slide_file(slide)

        if slide.original_filename.lower().endswith(".pdf"):
            try:
                import fitz
                doc = fitz.open(str(file_path))

                if page_number < 1 or page_number > len(doc):
                    doc.close()
                    raise HTTPException(
                        status_code=400,
                        detail=f"Invalid page number. File has {len(doc)} pages.",
                    )

                page = doc[page_number - 1]
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                doc.close()

                return Response(content=img_bytes, media_type="image/png")

            except ImportError:
                raise HTTPException(status_code=500, detail="PyMuPDF not installed for PDF rendering")
            except HTTPException:
                raise
            except Exception as e:
                logger.error(f"Error rendering PDF page: {e}")
                raise HTTPException(status_code=500, detail="Error rendering PDF")

        elif slide.original_filename.lower().endswith((".ppt", ".pptx")):
            try:
                with tempfile.TemporaryDirectory() as temp_dir:
                    try:
                        result = await asyncio.to_thread(
                            subprocess.run,
                            [
                                "soffice",
                                "--headless",
                                "--convert-to",
                                "pdf",
                                "--outdir",
                                temp_dir,
                                str(file_path),
                            ],
                            capture_output=True,
                            timeout=60,
                        )

                        converted_files = [f for f in os.listdir(temp_dir) if f.endswith(".pdf")]
                        if converted_files:
                            pdf_path = os.path.join(temp_dir, converted_files[0])

                            import fitz
                            doc = fitz.open(pdf_path)

                            if page_number < 1 or page_number > len(doc):
                                doc.close()
                                raise HTTPException(status_code=400, detail="Invalid page number")

                            page = doc[page_number - 1]
                            mat = fitz.Matrix(2.0, 2.0)
                            pix = page.get_pixmap(matrix=mat)
                            img_bytes = pix.tobytes("png")
                            doc.close()

                            return Response(content=img_bytes, media_type="image/png")
                    except (subprocess.TimeoutExpired, FileNotFoundError):
                        pass

                try:
                    from pptx import Presentation as PptxPresentation
                    from pptx.util import Inches, Pt
                    from PIL import Image, ImageDraw, ImageFont

                    prs = PptxPresentation(str(file_path))

                    if page_number < 1 or page_number > len(prs.slides):
                        raise HTTPException(
                            status_code=400,
                            detail=f"Invalid page number. File has {len(prs.slides)} slides.",
                        )

                    ppt_slide = prs.slides[page_number - 1]

                    width, height = 1280, 720
                    img = Image.new("RGB", (width, height), color="#ffffff")
                    draw = ImageDraw.Draw(img)

                    try:
                        title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 36)
                        body_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
                        small_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
                    except Exception:
                        try:
                            title_font = ImageFont.truetype(
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36
                            )
                            body_font = ImageFont.truetype(
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20
                            )
                            small_font = ImageFont.truetype(
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16
                            )
                        except Exception:
                            title_font = ImageFont.load_default()
                            body_font = ImageFont.load_default()
                            small_font = ImageFont.load_default()

                    y_position = 40
                    texts_extracted = []

                    for shape in ppt_slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts_extracted.append(shape.text.strip())

                    if texts_extracted:
                        title = texts_extracted[0][:100]
                        words = title.split()
                        lines = []
                        current_line = []
                        for word in words:
                            current_line.append(word)
                            test_line = " ".join(current_line)
                            bbox = draw.textbbox((0, 0), test_line, font=title_font)
                            if bbox[2] > width - 80:
                                current_line.pop()
                                if current_line:
                                    lines.append(" ".join(current_line))
                                current_line = [word]
                        if current_line:
                            lines.append(" ".join(current_line))

                        for line in lines[:2]:
                            draw.text((40, y_position), line, fill="#1a1a2e", font=title_font)
                            y_position += 45

                        y_position += 20

                    for text in texts_extracted[1:]:
                        if y_position > height - 60:
                            break

                        words = text.split()
                        lines = []
                        current_line = []
                        for word in words:
                            current_line.append(word)
                            test_line = " ".join(current_line)
                            bbox = draw.textbbox((0, 0), test_line, font=body_font)
                            if bbox[2] > width - 100:
                                current_line.pop()
                                if current_line:
                                    lines.append(" ".join(current_line))
                                current_line = [word]
                        if current_line:
                            lines.append(" ".join(current_line))

                        for line in lines[:4]:
                            if y_position > height - 60:
                                break
                            draw.text((50, y_position), f"• {line}", fill="#333333", font=body_font)
                            y_position += 28

                        y_position += 15

                    draw.text((width - 60, height - 30), f"{page_number}", fill="#888888", font=small_font)

                    if not texts_extracted:
                        draw.text(
                            (width // 2, height // 2 - 20),
                            f"Slide {page_number}",
                            fill="#1a1a2e",
                            font=title_font,
                            anchor="mm",
                        )
                        draw.text(
                            (width // 2, height // 2 + 30),
                            "(Content preview not available)",
                            fill="#888888",
                            font=body_font,
                            anchor="mm",
                        )

                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format="PNG", quality=95)
                    img_bytes = img_buffer.getvalue()

                    return Response(content=img_bytes, media_type="image/png")

                except ImportError as e:
                    logger.error(f"python-pptx not available: {e}")
                    from PIL import Image, ImageDraw

                    img = Image.new("RGB", (800, 600), color="#1a1a2e")
                    draw = ImageDraw.Draw(img)
                    draw.text((400, 280), f"Slide {page_number}", fill="#d7b38c", anchor="mm")
                    draw.text((400, 320), "Preview not available", fill="#888888", anchor="mm")

                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format="PNG")
                    img_bytes = img_buffer.getvalue()

                    return Response(content=img_bytes, media_type="image/png")

            except HTTPException:
                raise
            except Exception as e:
                logger.error(f"Error rendering PowerPoint: {e}")
                raise HTTPException(status_code=500, detail="Error rendering PowerPoint")

        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting slide image: {str(e)}")
        raise HTTPException(status_code=500, detail="Error getting slide image")

@router.get("/media/podcast/voice-modes")
def get_podcast_voice_modes():
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")
    return {"voice_modes": podcast_agent_service.list_voice_modes()}

@router.get("/media/podcast/voice-personas")
def get_podcast_voice_personas():
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")
    return {"voice_personas": podcast_agent_service.list_voice_personas()}

@router.get("/media/podcast/languages")
def get_podcast_languages():
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")
    return {"languages": podcast_agent_service.list_supported_languages()}

@router.get("/media/podcast/difficulties")
def get_podcast_difficulties():
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")
    return {"difficulties": podcast_agent_service.list_difficulty_levels()}

@router.post("/media/podcast/start")
async def start_media_podcast(
    payload: PodcastStartRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        session_data = podcast_agent_service.start_session(
            db=db,
            user_id=requested_user.id,
            transcript=payload.transcript,
            analysis=payload.analysis,
            title=payload.title,
            source_type=payload.source_type,
            voice_mode=payload.voice_mode,
            voice_persona=payload.voice_persona,
            difficulty=payload.difficulty,
            answer_language=payload.answer_language,
            session_options=payload.session_options,
        )
        return {"success": True, **session_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to start media podcast: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to start podcast session")

@router.post("/media/podcast/next")
async def next_media_podcast_segment(
    payload: PodcastSessionRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        segment_data = podcast_agent_service.get_next_segment(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
        )
        return {"success": True, **segment_data}
    except ValueError as e:
        raise HTTPException(status_code=404, detail="Not found")
    except Exception as e:
        logger.error("Failed to fetch next podcast segment: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to fetch next podcast segment")

@router.post("/media/podcast/ask")
async def ask_media_podcast_question(
    payload: PodcastQuestionRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        answer_data = podcast_agent_service.ask_question(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            question=payload.question,
            voice_mode=payload.voice_mode,
            voice_persona=payload.voice_persona,
            difficulty=payload.difficulty,
            question_language=payload.question_language,
            answer_language=payload.answer_language,
        )
        return {"success": True, **answer_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Podcast question handling failed: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to answer podcast question")

@router.post("/media/podcast/voice-mode")
async def set_media_podcast_voice_mode(
    payload: PodcastVoiceModeRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        mode_data = podcast_agent_service.set_voice_mode(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            voice_mode=payload.voice_mode,
        )
        return {"success": True, **mode_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to update podcast voice mode: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to update podcast voice mode")

@router.post("/media/podcast/settings")
async def update_media_podcast_settings(
    payload: PodcastSettingsRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        settings_data = podcast_agent_service.update_settings(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            difficulty=payload.difficulty,
            answer_language=payload.answer_language,
            voice_persona=payload.voice_persona,
            voice_mode=payload.voice_mode,
            session_options=payload.session_options,
        )
        return {"success": True, **settings_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to update podcast settings: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to update podcast settings")

@router.post("/media/podcast/bookmark")
async def add_media_podcast_bookmark(
    payload: PodcastBookmarkRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        bookmark = podcast_agent_service.add_bookmark(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            label=payload.label,
            chapter_index=payload.chapter_index,
            timestamp_seconds=payload.timestamp_seconds,
        )
        bookmarks = podcast_agent_service.list_bookmarks(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
        )
        return {"success": True, "bookmark": bookmark, "bookmarks": bookmarks}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to add podcast bookmark: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to add podcast bookmark")

@router.post("/media/podcast/bookmarks")
async def get_media_podcast_bookmarks(
    payload: PodcastSessionRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        bookmarks = podcast_agent_service.list_bookmarks(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
        )
        return {"success": True, "bookmarks": bookmarks}
    except Exception as e:
        logger.error("Failed to fetch podcast bookmarks: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to fetch podcast bookmarks")

@router.post("/media/podcast/replay")
async def replay_media_podcast_bookmark(
    payload: PodcastReplayBookmarkRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        replay_data = podcast_agent_service.replay_bookmark(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            bookmark_id=payload.bookmark_id,
        )
        return {"success": True, **replay_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to replay bookmark: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to replay bookmark")

@router.post("/media/podcast/jump")
async def jump_media_podcast_chapter(
    payload: PodcastJumpChapterRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        jump_data = podcast_agent_service.jump_to_chapter(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            chapter_index=payload.chapter_index,
        )
        return {"success": True, **jump_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to jump podcast chapter: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to jump podcast chapter")

@router.post("/media/podcast/mcq/start")
async def start_media_podcast_mcq(
    payload: PodcastMCQStartRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        mcq_data = podcast_agent_service.start_mcq_drill(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            count=payload.count,
        )
        return {"success": True, **mcq_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to start podcast mcq drill: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to start podcast mcq drill")

@router.post("/media/podcast/mcq/answer")
async def answer_media_podcast_mcq(
    payload: PodcastMCQAnswerRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        mcq_data = podcast_agent_service.answer_mcq(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            question_index=payload.question_index,
            selected_index=payload.selected_index,
        )
        return {"success": True, **mcq_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to answer podcast mcq: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to answer podcast mcq")

@router.post("/media/podcast/export")
async def export_media_podcast_session(
    payload: PodcastExportRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        export_data = podcast_agent_service.export_session(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
            format_type=payload.format_type,
        )
        return {"success": True, **export_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to export podcast session: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to export podcast session")

@router.get("/media/podcast/sessions")
async def list_media_podcast_sessions(
    user_id: str = Query(...),
    limit: int = Query(20),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, user_id)
    _verify_user_access(requested_user, current_user)

    try:
        sessions = podcast_agent_service.list_user_sessions(
            db=db,
            user_id=requested_user.id,
            limit=limit,
        )
        return {"success": True, "sessions": sessions}
    except Exception as e:
        logger.error("Failed to list podcast sessions: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to list podcast sessions")

@router.post("/media/podcast/resume")
async def resume_media_podcast_session(
    payload: PodcastSessionRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        resume_data = podcast_agent_service.resume_session(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
        )
        return {"success": True, **resume_data}
    except ValueError as e:
        raise HTTPException(status_code=404, detail="Not found")
    except Exception as e:
        logger.error("Failed to resume podcast session: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to resume podcast session")

@router.post("/media/podcast/stop")
async def stop_media_podcast(
    payload: PodcastSessionRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        stop_data = podcast_agent_service.stop_session(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
        )
        return {"success": True, **stop_data}
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Bad request")
    except Exception as e:
        logger.error("Failed to stop podcast session: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to stop podcast session")

@router.post("/media/podcast/state")
async def get_media_podcast_state(
    payload: PodcastSessionRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not podcast_agent_service:
        raise HTTPException(status_code=500, detail="Podcast agent is not available")

    requested_user = _resolve_user_or_404(db, payload.user_id)
    _verify_user_access(requested_user, current_user)

    try:
        state_data = podcast_agent_service.get_state(
            db=db,
            session_id=payload.session_id,
            user_id=requested_user.id,
        )
        return {"success": True, **state_data}
    except ValueError as e:
        raise HTTPException(status_code=404, detail="Not found")
    except Exception as e:
        logger.error("Failed to fetch podcast session state: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to fetch podcast session state")
